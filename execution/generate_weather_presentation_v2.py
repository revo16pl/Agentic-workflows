#!/usr/bin/env python3
"""
Weather Presentation V2 - Professional Design
Matches ABB template with gradients, charts, and proper spacing
"""

import os
import sys
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import random
import io

class WeatherPresentationV2:
    def __init__(self, output_dir='test_02'):
        self.output_dir = output_dir
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

        # Professional color palette (from UX/UI analysis)
        self.primary = RGBColor(30, 64, 175)      # #1E40AF
        self.secondary = RGBColor(59, 130, 246)    # #3B82F6
        self.accent = RGBColor(245, 158, 11)       # #F59E0B
        self.success = RGBColor(16, 185, 129)      # #10B981
        self.warning = RGBColor(239, 68, 68)       # #EF4444
        self.text_dark = RGBColor(30, 58, 138)     # #1E3A8A
        self.bg_light = RGBColor(248, 250, 252)    # #F8FAFC
        self.white = RGBColor(255, 255, 255)
        self.gray = RGBColor(100, 116, 139)        # #64748B
        self.light_gray = RGBColor(226, 232, 240)  # #E2E8F0

        # ABB Red for accent elements
        self.red = RGBColor(255, 0, 0)

        # Typography scale
        self.font_family = 'Arial'

        # Generate weather data
        self.weather_data = self.generate_fake_data()

    def generate_fake_data(self):
        """Generate realistic fake weather data"""
        base_date = datetime.now()

        data = {
            'city': 'Warszawa',
            'current': {
                'temp': 8.5,
                'feels_like': 6.2,
                'humidity': 72,
                'wind_speed': 15,
                'pressure': 1013,
                'condition': 'Pochmurnie z przejaśnieniami'
            },
            'daily': [],
            'stats': {
                'avg_temp': 7.2,
                'max_temp': 12.5,
                'min_temp': 2.1,
                'rainy_days': 3,
                'sunny_days': 2,
                'cloudy_days': 2
            },
            'hourly_temps': []
        }

        # Daily forecast
        conditions = ['Słonecznie', 'Pochmurnie', 'Deszczowo', 'Częściowo pochmurnie']
        icons = ['☀️', '☁️', '🌧️', '⛅']

        for i in range(7):
            day = base_date + timedelta(days=i)
            cond_idx = random.randint(0, 3)
            data['daily'].append({
                'date': day.strftime('%d.%m'),
                'day': ['Pn', 'Wt', 'Śr', 'Cz', 'Pt', 'So', 'Nd'][day.weekday()],
                'temp_max': random.randint(8, 14),
                'temp_min': random.randint(2, 7),
                'condition': conditions[cond_idx],
                'icon': icons[cond_idx],
                'rain': random.randint(0, 60),
                'wind': random.randint(5, 25)
            })

        # Hourly temperatures for chart
        for hour in range(24):
            temp = 5 + 3 * (1 - abs(hour - 14) / 14) + random.uniform(-0.5, 0.5)
            data['hourly_temps'].append({
                'hour': hour,
                'temp': round(temp, 1)
            })

        return data

    def create_gradient_background(self, width=1000, height=750):
        """Create red-pink-blue gradient background image"""
        img = Image.new('RGB', (width, height))
        draw = ImageDraw.Draw(img)

        # Define gradient colors: red -> pink -> blue
        colors = [
            (255, 0, 0),      # Red top
            (255, 100, 150),  # Pink middle
            (150, 100, 255)   # Blue bottom
        ]

        # Create vertical gradient
        for y in range(height):
            progress = y / height

            if progress < 0.5:
                # Red to Pink
                ratio = progress * 2
                r = int(colors[0][0] + (colors[1][0] - colors[0][0]) * ratio)
                g = int(colors[0][1] + (colors[1][1] - colors[0][1]) * ratio)
                b = int(colors[0][2] + (colors[1][2] - colors[0][2]) * ratio)
            else:
                # Pink to Blue
                ratio = (progress - 0.5) * 2
                r = int(colors[1][0] + (colors[2][0] - colors[1][0]) * ratio)
                g = int(colors[1][1] + (colors[2][1] - colors[1][1]) * ratio)
                b = int(colors[1][2] + (colors[2][2] - colors[1][2]) * ratio)

            draw.line([(0, y), (width, y)], fill=(r, g, b))

        # Save to bytes
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        return img_bytes

    def create_temperature_chart(self):
        """Create line chart for 7-day temperature forecast"""
        fig, ax = plt.subplots(figsize=(8, 4), facecolor='white')

        days = [d['day'] + '\n' + d['date'] for d in self.weather_data['daily']]
        temp_max = [d['temp_max'] for d in self.weather_data['daily']]
        temp_min = [d['temp_min'] for d in self.weather_data['daily']]

        # Plot lines
        ax.plot(days, temp_max, marker='o', linewidth=3, markersize=8,
                color='#EF4444', label='Maksymalna')
        ax.plot(days, temp_min, marker='o', linewidth=3, markersize=8,
                color='#3B82F6', label='Minimalna')

        # Fill between
        ax.fill_between(range(len(days)), temp_min, temp_max, alpha=0.2, color='#64748B')

        # Styling
        ax.set_ylabel('Temperatura (°C)', fontsize=12, fontweight='bold')
        ax.legend(loc='upper right', frameon=False, fontsize=10)
        ax.grid(True, alpha=0.3, linestyle='--')
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)

        plt.tight_layout()

        # Save to bytes
        img_bytes = io.BytesIO()
        plt.savefig(img_bytes, format='PNG', dpi=150, bbox_inches='tight')
        img_bytes.seek(0)
        plt.close()

        return img_bytes

    def create_hourly_temp_chart(self):
        """Create area chart for hourly temperatures"""
        fig, ax = plt.subplots(figsize=(8, 3), facecolor='white')

        hours = [h['hour'] for h in self.weather_data['hourly_temps']]
        temps = [h['temp'] for h in self.weather_data['hourly_temps']]

        # Area chart
        ax.fill_between(hours, temps, alpha=0.3, color='#3B82F6')
        ax.plot(hours, temps, linewidth=2, color='#1E40AF')

        # Styling
        ax.set_xlabel('Godzina', fontsize=10)
        ax.set_ylabel('°C', fontsize=10)
        ax.set_xticks([0, 6, 12, 18, 23])
        ax.set_xticklabels(['00:00', '06:00', '12:00', '18:00', '23:00'])
        ax.grid(True, alpha=0.3, axis='y')
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)

        plt.tight_layout()

        # Save to bytes
        img_bytes = io.BytesIO()
        plt.savefig(img_bytes, format='PNG', dpi=150, bbox_inches='tight')
        img_bytes.seek(0)
        plt.close()

        return img_bytes

    def create_stats_comparison(self):
        """Create bar chart for temperature statistics"""
        fig, ax = plt.subplots(figsize=(6, 3), facecolor='white')

        stats = self.weather_data['stats']
        categories = ['Średnia', 'Maks', 'Min']
        values = [stats['avg_temp'], stats['max_temp'], stats['min_temp']]
        colors = ['#F59E0B', '#EF4444', '#3B82F6']

        bars = ax.barh(categories, values, color=colors, height=0.6)

        # Add value labels
        for i, (bar, val) in enumerate(zip(bars, values)):
            ax.text(val + 0.3, i, f'{val}°C', va='center', fontsize=11, fontweight='bold')

        ax.set_xlabel('Temperatura (°C)', fontsize=10)
        ax.set_xlim(0, max(values) + 2)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)

        plt.tight_layout()

        # Save to bytes
        img_bytes = io.BytesIO()
        plt.savefig(img_bytes, format='PNG', dpi=150, bbox_inches='tight')
        img_bytes.seek(0)
        plt.close()

        return img_bytes

    def add_blank_slide(self):
        """Add blank slide"""
        blank_layout = self.prs.slide_layouts[6]
        return self.prs.slides.add_slide(blank_layout)

    def add_footer(self, slide, slide_num):
        """Add professional footer"""
        # Copyright
        footer_left = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(3), Inches(0.2))
        tf = footer_left.text_frame
        tf.text = f"© {datetime.now().year} Weather Service"
        tf.paragraphs[0].font.size = Pt(8)
        tf.paragraphs[0].font.color.rgb = self.gray
        tf.paragraphs[0].font.name = self.font_family

        # Slide number
        footer_center = slide.shapes.add_textbox(Inches(4.5), Inches(7.2), Inches(1), Inches(0.2))
        tf = footer_center.text_frame
        tf.text = f"{slide_num}"
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = self.gray
        tf.paragraphs[0].font.name = self.font_family
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Date
        footer_right = slide.shapes.add_textbox(Inches(7.5), Inches(7.2), Inches(2), Inches(0.2))
        tf = footer_right.text_frame
        tf.text = f"Luty {datetime.now().year}"
        tf.paragraphs[0].font.size = Pt(8)
        tf.paragraphs[0].font.color.rgb = self.gray
        tf.paragraphs[0].font.name = self.font_family
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def add_title_slide(self):
        """Title slide with gradient background"""
        slide = self.add_blank_slide()

        # Add gradient background image
        gradient_img = self.create_gradient_background()
        left = Inches(0)
        top = Inches(0)
        pic = slide.shapes.add_picture(gradient_img, left, top,
                                      width=self.prs.slide_width,
                                      height=self.prs.slide_height)
        # Send to back
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        # Main title - large and bold
        title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
        tf = title.text_frame
        tf.text = "Prognoza Pogody\nWarszawa"
        tf.word_wrap = True

        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(72)
            paragraph.font.color.rgb = self.white
            paragraph.font.bold = True
            paragraph.font.name = self.font_family

        # Subtitle
        subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(6), Inches(0.6))
        tf = subtitle.text_frame
        tf.text = "LUTY 2026 • WARSZAWA"
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = self.white
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = self.font_family

        # Author
        author = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(6), Inches(0.4))
        tf = author.text_frame
        tf.text = "Raport 7-dniowy • Analiza temperatury i opadów"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].font.name = self.font_family

        # Logo bottom right
        logo = slide.shapes.add_textbox(Inches(7.5), Inches(6.8), Inches(2), Inches(0.5))
        tf = logo.text_frame
        tf.text = "WEATHER\nSERVICE"
        for p in tf.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = self.white
            p.font.bold = True
            p.font.name = self.font_family
            p.alignment = PP_ALIGN.RIGHT

        # Copyright
        copy = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(4), Inches(0.2))
        tf = copy.text_frame
        tf.text = f"© {datetime.now().year} Weather Service. Wszelkie prawa zastrzeżone."
        tf.paragraphs[0].font.size = Pt(8)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].font.name = self.font_family

    def add_toc_slide(self):
        """Table of contents with clean layout"""
        slide = self.add_blank_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.white

        # Red accent line
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(0.6), Inches(0.15), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = self.red
        line.line.fill.background()

        # Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(8), Inches(0.8))
        tf = title.text_frame
        tf.text = "Spis treści"
        tf.paragraphs[0].font.size = Pt(42)
        tf.paragraphs[0].font.color.rgb = self.text_dark
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = self.font_family

        # TOC items with better spacing
        items = [
            "I.     Kontekst i dane podstawowe",
            "II.    Prognoza 7-dniowa",
            "III.   Analiza temperatur i trendy",
            "IV.   Opady i wiatr",
            "V.    Podsumowanie i rekomendacje"
        ]

        content = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(8), Inches(4))
        tf = content.text_frame

        for i, item in enumerate(items):
            if i == 0:
                tf.text = item
            else:
                p = tf.add_paragraph()
                p.text = item
            tf.paragraphs[i].font.size = Pt(22)
            tf.paragraphs[i].font.color.rgb = self.text_dark
            tf.paragraphs[i].font.name = self.font_family
            tf.paragraphs[i].space_after = Pt(24)

        self.add_footer(slide, 2)

    def add_section_divider(self, num, title):
        """Section divider with gradient"""
        slide = self.add_blank_slide()

        # Add gradient background
        gradient_img = self.create_gradient_background()
        pic = slide.shapes.add_picture(gradient_img, Inches(0), Inches(0),
                                      width=self.prs.slide_width,
                                      height=self.prs.slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        # Large section number
        number = slide.shapes.add_textbox(Inches(7), Inches(1), Inches(2.5), Inches(2))
        tf = number.text_frame
        tf.text = f"0{num}"
        tf.paragraphs[0].font.size = Pt(140)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = self.font_family
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Section title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(7), Inches(2))
        tf = title_box.text_frame
        tf.text = title
        tf.word_wrap = True
        tf.paragraphs[0].font.size = Pt(66)
        tf.paragraphs[0].font.color.rgb = self.white
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = self.font_family

    def add_content_slide_with_header(self, title, slide_num):
        """Content slide with red accent line"""
        slide = self.add_blank_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.bg_light

        # Red accent line
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(0.6), Inches(0.15), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = self.red
        line.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(8), Inches(0.6))
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.color.rgb = self.text_dark
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = self.font_family

        self.add_footer(slide, slide_num)
        return slide

    def add_current_weather_slide(self):
        """Current weather with icons and metrics"""
        slide = self.add_content_slide_with_header("Aktualne warunki pogodowe", 4)

        current = self.weather_data['current']

        # Large temperature display
        temp_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(2))
        tf = temp_box.text_frame

        # Temperature
        p = tf.paragraphs[0]
        p.text = f"{current['temp']}°C"
        p.font.size = Pt(84)
        p.font.color.rgb = self.primary
        p.font.bold = True
        p.font.name = self.font_family

        # Condition
        p = tf.add_paragraph()
        p.text = current['condition']
        p.font.size = Pt(18)
        p.font.color.rgb = self.text_dark
        p.font.name = self.font_family
        p.space_before = Pt(12)

        # Feels like
        p = tf.add_paragraph()
        p.text = f"Odczuwalna: {current['feels_like']}°C"
        p.font.size = Pt(14)
        p.font.color.rgb = self.gray
        p.font.name = self.font_family
        p.space_before = Pt(8)

        # KPI Cards on the right
        metrics = [
            ('Wilgotność', f"{current['humidity']}%", self.secondary),
            ('Wiatr', f"{current['wind_speed']} km/h", self.accent),
            ('Ciśnienie', f"{current['pressure']} hPa", self.success)
        ]

        y_pos = 2
        for label, value, color in metrics:
            # Card background
            card = slide.shapes.add_shape(1, Inches(5.5), Inches(y_pos), Inches(3.8), Inches(1.2))
            card.fill.solid()
            card.fill.fore_color.rgb = self.white
            card.line.color.rgb = self.light_gray
            card.line.width = Pt(1)

            # Label
            lbl = slide.shapes.add_textbox(Inches(5.7), Inches(y_pos + 0.15), Inches(3.4), Inches(0.3))
            lbl.text_frame.text = label
            lbl.text_frame.paragraphs[0].font.size = Pt(12)
            lbl.text_frame.paragraphs[0].font.color.rgb = self.gray
            lbl.text_frame.paragraphs[0].font.name = self.font_family

            # Value
            val = slide.shapes.add_textbox(Inches(5.7), Inches(y_pos + 0.5), Inches(3.4), Inches(0.6))
            val.text_frame.text = value
            val.text_frame.paragraphs[0].font.size = Pt(32)
            val.text_frame.paragraphs[0].font.color.rgb = color
            val.text_frame.paragraphs[0].font.bold = True
            val.text_frame.paragraphs[0].font.name = self.font_family

            y_pos += 1.4

    def add_weekly_forecast_slide(self):
        """7-day forecast with chart"""
        slide = self.add_content_slide_with_header("Prognoza 7-dniowa", 7)

        # Add temperature chart
        chart_img = self.create_temperature_chart()
        pic = slide.shapes.add_picture(chart_img, Inches(0.8), Inches(2), width=Inches(8.4))

        # Add daily cards below chart
        daily = self.weather_data['daily']
        x_start = 0.5
        col_width = 1.35

        y_base = 5.5

        for i, day in enumerate(daily):
            x_pos = x_start + (i * col_width)

            # Day card
            card = slide.shapes.add_shape(1, Inches(x_pos), Inches(y_base), 
                                         Inches(col_width - 0.05), Inches(1.3))
            card.fill.solid()
            card.fill.fore_color.rgb = self.white
            card.line.color.rgb = self.light_gray
            card.line.width = Pt(0.5)

            # Day name
            day_name = slide.shapes.add_textbox(Inches(x_pos + 0.05), Inches(y_base + 0.1),
                                               Inches(col_width - 0.15), Inches(0.25))
            tf = day_name.text_frame
            tf.text = f"{day['day']}"
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = self.text_dark
            tf.paragraphs[0].font.name = self.font_family
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Icon (emoji as placeholder)
            icon = slide.shapes.add_textbox(Inches(x_pos + 0.05), Inches(y_base + 0.4),
                                           Inches(col_width - 0.15), Inches(0.3))
            tf = icon.text_frame
            tf.text = day['icon']
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Temp high/low
            temps = slide.shapes.add_textbox(Inches(x_pos + 0.05), Inches(y_base + 0.75),
                                            Inches(col_width - 0.15), Inches(0.4))
            tf = temps.text_frame
            tf.text = f"{day['temp_max']}° / {day['temp_min']}°"
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = self.text_dark
            tf.paragraphs[0].font.name = self.font_family
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_hourly_temp_slide(self):
        """Hourly temperature analysis"""
        slide = self.add_content_slide_with_header("Analiza temperatury godzinowej", 9)

        # Add hourly chart
        chart_img = self.create_hourly_temp_chart()
        pic = slide.shapes.add_picture(chart_img, Inches(0.8), Inches(2.2), width=Inches(8.4))

        # Add key insights
        insights = [
            f"🌅 Najniższa temperatura: {min(h['temp'] for h in self.weather_data['hourly_temps']):.1f}°C (wczesnym rankiem)",
            f"☀️ Najwyższa temperatura: {max(h['temp'] for h in self.weather_data['hourly_temps']):.1f}°C (po południu)",
            f"📊 Średnia dzienna: {sum(h['temp'] for h in self.weather_data['hourly_temps'])/24:.1f}°C"
        ]

        y_pos = 5.5
        for insight in insights:
            txt = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.3))
            tf = txt.text_frame
            tf.text = insight
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = self.text_dark
            tf.paragraphs[0].font.name = self.font_family
            y_pos += 0.35

    def add_stats_slide(self):
        """Statistics with bar chart"""
        slide = self.add_content_slide_with_header("Statystyki tygodniowe", 12)

        stats = self.weather_data['stats']

        # Add stats chart
        chart_img = self.create_stats_comparison()
        pic = slide.shapes.add_picture(chart_img, Inches(1.5), Inches(2.3), width=Inches(7))

        # Weather distribution
        weather_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
        tf = weather_box.text_frame

        p = tf.paragraphs[0]
        p.text = "Rozkład pogody w tygodniu:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.text_dark
        p.font.name = self.font_family

        p = tf.add_paragraph()
        p.text = f"☀️ Słonecznie: {stats['sunny_days']} dni  •  ☁️ Pochmurnie: {stats['cloudy_days']} dni  •  🌧️ Deszczowo: {stats['rainy_days']} dni"
        p.font.size = Pt(14)
        p.font.color.rgb = self.text_dark
        p.font.name = self.font_family
        p.space_before = Pt(12)

    def add_summary_slide(self):
        """Summary with key takeaways"""
        slide = self.add_content_slide_with_header("Podsumowanie i rekomendacje", 14)

        stats = self.weather_data['stats']

        # Summary cards
        summaries = [
            {
                'icon': '✅',
                'title': 'Stabilne warunki pogodowe',
                'text': f"Średnia temperatura {stats['avg_temp']}°C z wahaniami od {stats['min_temp']}°C do {stats['max_temp']}°C. Warunki sprzyjają planowaniu aktywności."
            },
            {
                'icon': '☔',
                'title': 'Umiarkowane opady',
                'text': f"Przewidywane {stats['rainy_days']} dni z opadami. Zalecane noszenie parasolki w dni robocze, szczególnie w godzinach popołudniowych."
            },
            {
                'icon': '🌡️',
                'title': 'Komfort termiczny',
                'text': f"Temperatura odczuwalna około {self.weather_data['current']['feels_like']}°C. Wilgotność {self.weather_data['current']['humidity']}% zapewnia komfortowe warunki."
            }
        ]

        y_pos = 2.2
        for summary in summaries:
            # Card
            card = slide.shapes.add_shape(1, Inches(1), Inches(y_pos), Inches(8), Inches(1.3))
            card.fill.solid()
            card.fill.fore_color.rgb = self.white
            card.line.color.rgb = self.light_gray
            card.line.width = Pt(1)

            # Icon
            icon_box = slide.shapes.add_textbox(Inches(1.3), Inches(y_pos + 0.15), Inches(0.5), Inches(0.5))
            tf = icon_box.text_frame
            tf.text = summary['icon']
            tf.paragraphs[0].font.size = Pt(36)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Title
            title_box = slide.shapes.add_textbox(Inches(2.2), Inches(y_pos + 0.2), Inches(6.5), Inches(0.35))
            tf = title_box.text_frame
            tf.text = summary['title']
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = self.text_dark
            tf.paragraphs[0].font.name = self.font_family

            # Text
            text_box = slide.shapes.add_textbox(Inches(2.2), Inches(y_pos + 0.6), Inches(6.5), Inches(0.6))
            tf = text_box.text_frame
            tf.text = summary['text']
            tf.word_wrap = True
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.color.rgb = self.text_dark
            tf.paragraphs[0].font.name = self.font_family

            y_pos += 1.5

    def add_final_slide(self):
        """Final slide with logo"""
        slide = self.add_blank_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.bg_light

        # Large logo/text
        logo = slide.shapes.add_textbox(Inches(2.5), Inches(2.8), Inches(5), Inches(2))
        tf = logo.text_frame
        tf.text = "WEATHER\nSERVICE"
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        for p in tf.paragraphs:
            p.font.size = Pt(68)
            p.font.color.rgb = self.primary
            p.font.bold = True
            p.font.name = self.font_family
            p.alignment = PP_ALIGN.CENTER

        # Tagline
        tagline = slide.shapes.add_textbox(Inches(2.5), Inches(5), Inches(5), Inches(0.5))
        tf = tagline.text_frame
        tf.text = "Twój partner w prognozowaniu pogody"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = self.gray
        tf.paragraphs[0].font.name = self.font_family
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def generate(self):
        """Generate complete presentation"""
        print("\n" + "="*60)
        print("  WEATHER PRESENTATION V2 - PROFESSIONAL DESIGN")
        print("="*60 + "\n")

        print("📊 Generowanie prezentacji z wykresami i gradientami...\n")

        # Slide 1: Title with gradient
        print("  ✓ [1/15] Slajd tytułowy z gradientem")
        self.add_title_slide()

        # Slide 2: TOC
        print("  ✓ [2/15] Spis treści")
        self.add_toc_slide()

        # Slide 3: Section 1
        print("  ✓ [3/15] Divider: Kontekst i dane")
        self.add_section_divider(1, "Kontekst i dane\npodstawowe")

        # Slide 4: Current weather
        print("  ✓ [4/15] Aktualne warunki pogodowe")
        self.add_current_weather_slide()

        # Slide 5: Section 2
        print("  ✓ [5/15] Divider: Prognoza tygodniowa")
        self.add_section_divider(2, "Prognoza\ntygodniowa")

        # Slide 6: Placeholder
        print("  ✓ [6/15] Placeholder slide")
        slide = self.add_content_slide_with_header("Trend temperatury", 6)

        # Slide 7: Weekly forecast with chart
        print("  ✓ [7/15] Prognoza 7-dniowa z wykresem")
        self.add_weekly_forecast_slide()

        # Slide 8: Section 3
        print("  ✓ [8/15] Divider: Analiza temperatur")
        self.add_section_divider(3, "Analiza\ntemperatur")

        # Slide 9: Hourly temps
        print("  ✓ [9/15] Analiza temperatury godzinowej")
        self.add_hourly_temp_slide()

        # Slide 10: Section 4
        print("  ✓ [10/15] Divider: Opady i wiatr")
        self.add_section_divider(4, "Opady\ni wiatr")

        # Slide 11: Placeholder
        print("  ✓ [11/15] Placeholder slide")
        slide = self.add_content_slide_with_header("Analiza opadów", 11)

        # Slide 12: Statistics
        print("  ✓ [12/15] Statystyki z wykresem")
        self.add_stats_slide()

        # Slide 13: Section 5
        print("  ✓ [13/15] Divider: Podsumowanie")
        self.add_section_divider(5, "Podsumowanie")

        # Slide 14: Summary
        print("  ✓ [14/15] Podsumowanie i rekomendacje")
        self.add_summary_slide()

        # Slide 15: Final
        print("  ✓ [15/15] Slajd końcowy")
        self.add_final_slide()

        # Save
        output_path = os.path.join(self.output_dir, 'Prognoza_Pogody_Warszawa_v2.pptx')
        os.makedirs(self.output_dir, exist_ok=True)
        self.prs.save(output_path)

        print("\n" + "="*60)
        print(f"✅ Prezentacja zapisana: {output_path}")
        print("="*60 + "\n")

        return output_path


def main():
    try:
        gen = WeatherPresentationV2(output_dir='test_02')
        output_file = gen.generate()
        print(f"\n🎉 Sukces! Otwórz plik:\n   {output_file}\n")
        return 0
    except Exception as e:
        print(f"\n❌ Błąd: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        return 1


if __name__ == '__main__':
    sys.exit(main())
