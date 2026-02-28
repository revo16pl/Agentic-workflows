#!/usr/bin/env python3
"""
Generate Weather Presentation matching ABB template design
Creates PowerPoint with fake weather data for Warsaw
"""

import os
import sys
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import random

class WeatherPresentation:
    def __init__(self, output_dir='test_02'):
        self.output_dir = output_dir
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # ABB Colors
        self.red = RGBColor(255, 0, 0)
        self.white = RGBColor(255, 255, 255)
        self.black = RGBColor(0, 0, 0)
        self.gray = RGBColor(128, 128, 128)
        self.light_gray = RGBColor(240, 240, 240)
        
        self.weather_data = self.generate_fake_data()
    
    def generate_fake_data(self):
        """Generate fake weather data"""
        base_date = datetime.now()
        data = {
            'city': 'Warszawa',
            'current': {
                'temp': 8.5,
                'humidity': 72,
                'wind_speed': 15,
                'pressure': 1013
            },
            'daily': [],
            'stats': {
                'avg_temp': 7.2,
                'max_temp': 12.5,
                'min_temp': 2.1,
                'rainy_days': 3
            }
        }
        
        conditions = ['Słonecznie', 'Pochmurnie', 'Deszczowo', 'Mgliście']
        for i in range(7):
            day = base_date + timedelta(days=i)
            data['daily'].append({
                'date': day.strftime('%d.%m'),
                'day': ['Pn', 'Wt', 'Śr', 'Cz', 'Pt', 'So', 'Nd'][day.weekday()],
                'temp_max': random.randint(8, 14),
                'temp_min': random.randint(2, 7),
                'condition': random.choice(conditions),
                'rain': random.randint(0, 60)
            })
        
        return data
    
    def add_blank_slide(self):
        """Add blank slide"""
        blank_layout = self.prs.slide_layouts[6]
        return self.prs.slides.add_slide(blank_layout)
    
    def add_footer(self, slide, slide_num):
        """Add footer"""
        footer_left = slide.shapes.add_textbox(Inches(0.3), Inches(7.2), Inches(3), Inches(0.2))
        tf = footer_left.text_frame
        tf.text = f"© {datetime.now().year} Weather Service. All rights reserved."
        tf.paragraphs[0].font.size = Pt(8)
        tf.paragraphs[0].font.color.rgb = self.gray
        
        footer_center = slide.shapes.add_textbox(Inches(4), Inches(7.2), Inches(1), Inches(0.2))
        tf = footer_center.text_frame
        tf.text = f"Slide {slide_num}"
        tf.paragraphs[0].font.size = Pt(8)
        tf.paragraphs[0].font.color.rgb = self.gray
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        footer_right = slide.shapes.add_textbox(Inches(7.5), Inches(7.2), Inches(2), Inches(0.2))
        tf = footer_right.text_frame
        tf.text = f"Luty, {datetime.now().year}"
        tf.paragraphs[0].font.size = Pt(8)
        tf.paragraphs[0].font.color.rgb = self.gray
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    def add_title_slide(self):
        """Title slide with red background"""
        slide = self.add_blank_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.red
        
        # INTERNAL badge
        badge = slide.shapes.add_shape(1, Inches(0.3), Inches(0.3), Inches(1.2), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(100, 100, 100)
        badge.line.color.rgb = self.white
        tf = badge.text_frame
        tf.text = "INTERNAL"
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = self.white
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Main title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
        tf = title.text_frame
        tf.text = "Prognoza Pogody -\nWarszawa"
        for p in tf.paragraphs:
            p.font.size = Pt(66)
            p.font.color.rgb = self.white
            p.font.bold = True
        
        # Subtitle
        subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(6), Inches(1))
        tf = subtitle.text_frame
        tf.text = "LUTY 2026"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = self.white
        tf.paragraphs[0].font.bold = True
        
        # Author
        author = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(6), Inches(0.5))
        tf = author.text_frame
        tf.text = "DZIAŁ METEOROLOGII - ANALIZA DANYCH"
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = self.white
        
        # Tagline
        tagline = slide.shapes.add_textbox(Inches(7), Inches(6.8), Inches(2.5), Inches(0.5))
        tf = tagline.text_frame
        tf.text = "ENGINEERED\nTO OUTRUN"
        for p in tf.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = self.white
            p.font.bold = True
            p.alignment = PP_ALIGN.RIGHT
        
        # Copyright
        copy = slide.shapes.add_textbox(Inches(0.3), Inches(7.2), Inches(4), Inches(0.2))
        tf = copy.text_frame
        tf.text = f"© {datetime.now().year} Weather Service. All rights reserved."
        tf.paragraphs[0].font.size = Pt(8)
        tf.paragraphs[0].font.color.rgb = self.white
    
    def add_toc_slide(self):
        """Table of contents"""
        slide = self.add_blank_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.white
        
        # Red line
        line = slide.shapes.add_shape(1, Inches(0), Inches(0.5), Inches(0.3), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = self.red
        line.line.fill.background()
        
        # Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(8), Inches(0.6))
        tf = title.text_frame
        tf.text = "Spis treści"
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.color.rgb = self.black
        tf.paragraphs[0].font.bold = True
        
        # Items
        items = [
            "I.    Kontekst i dane podstawowe",
            "II.   Prognoza tygodniowa",
            "III.  Analiza temperatur",
            "IV.   Opady i wiatr",
            "V.    Podsumowanie"
        ]
        
        content = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8), Inches(4))
        tf = content.text_frame
        for i, item in enumerate(items):
            if i == 0:
                tf.text = item
            else:
                p = tf.add_paragraph()
                p.text = item
            tf.paragraphs[i].font.size = Pt(20)
            tf.paragraphs[i].font.color.rgb = self.black
            tf.paragraphs[i].space_after = Pt(20)
        
        self.add_footer(slide, 2)
    
    def add_section_divider(self, num, title):
        """Red section divider"""
        slide = self.add_blank_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.red
        
        # Number
        number = slide.shapes.add_textbox(Inches(7), Inches(0.8), Inches(2.5), Inches(1.5))
        tf = number.text_frame
        tf.text = f"0{num}"
        tf.paragraphs[0].font.size = Pt(120)
        tf.paragraphs[0].font.color.rgb = self.white
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(7), Inches(1.5))
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(66)
        tf.paragraphs[0].font.color.rgb = self.white
        tf.paragraphs[0].font.bold = True
    
    def add_content_slide(self, title, slide_num):
        """White content slide"""
        slide = self.add_blank_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.white
        
        # Red line
        line = slide.shapes.add_shape(1, Inches(0), Inches(0.5), Inches(0.3), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = self.red
        line.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(8), Inches(0.6))
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.color.rgb = self.black
        tf.paragraphs[0].font.bold = True
        
        self.add_footer(slide, slide_num)
        return slide
    
    def add_current_weather(self):
        """Current weather slide"""
        slide = self.add_content_slide("Aktualne warunki pogodowe – Warszawa", 4)
        current = self.weather_data['current']
        
        # Large temp
        temp = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(1.5))
        tf = temp.text_frame
        tf.text = f"{current['temp']}°C"
        tf.paragraphs[0].font.size = Pt(72)
        tf.paragraphs[0].font.color.rgb = self.red
        tf.paragraphs[0].font.bold = True
        
        # Metrics
        metrics = [
            ('Wilgotność:', f"{current['humidity']}%"),
            ('Wiatr:', f"{current['wind_speed']} km/h"),
            ('Ciśnienie:', f"{current['pressure']} hPa")
        ]
        
        y = 2.5
        for label, value in metrics:
            lbl = slide.shapes.add_textbox(Inches(5.5), Inches(y), Inches(2), Inches(0.3))
            lbl.text_frame.text = label
            lbl.text_frame.paragraphs[0].font.size = Pt(14)
            lbl.text_frame.paragraphs[0].font.color.rgb = self.gray
            
            val = slide.shapes.add_textbox(Inches(5.5), Inches(y+0.3), Inches(3), Inches(0.4))
            val.text_frame.text = value
            val.text_frame.paragraphs[0].font.size = Pt(24)
            val.text_frame.paragraphs[0].font.color.rgb = self.black
            val.text_frame.paragraphs[0].font.bold = True
            
            y += 1
    
    def add_weekly_forecast(self):
        """7-day forecast"""
        slide = self.add_content_slide("Prognoza 7-dniowa", 7)
        daily = self.weather_data['daily']
        
        x_start = 0.5
        col_width = 1.3
        
        for i, day in enumerate(daily):
            x = x_start + (i * col_width)
            
            # Day
            d = slide.shapes.add_textbox(Inches(x), Inches(2), Inches(col_width-0.1), Inches(0.4))
            d.text_frame.text = f"{day['day']}\n{day['date']}"
            d.text_frame.paragraphs[0].font.size = Pt(10)
            d.text_frame.paragraphs[0].font.bold = True
            d.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Condition
            c = slide.shapes.add_textbox(Inches(x), Inches(2.6), Inches(col_width-0.1), Inches(0.5))
            c.text_frame.text = day['condition'][:10]
            c.text_frame.paragraphs[0].font.size = Pt(8)
            c.text_frame.paragraphs[0].font.color.rgb = self.gray
            c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Temp high
            th = slide.shapes.add_textbox(Inches(x), Inches(3.3), Inches(col_width-0.1), Inches(0.4))
            th.text_frame.text = f"{day['temp_max']}°"
            th.text_frame.paragraphs[0].font.size = Pt(20)
            th.text_frame.paragraphs[0].font.color.rgb = self.red
            th.text_frame.paragraphs[0].font.bold = True
            th.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Temp low
            tl = slide.shapes.add_textbox(Inches(x), Inches(3.8), Inches(col_width-0.1), Inches(0.3))
            tl.text_frame.text = f"{day['temp_min']}°"
            tl.text_frame.paragraphs[0].font.size = Pt(12)
            tl.text_frame.paragraphs[0].font.color.rgb = self.gray
            tl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def add_stats(self):
        """Statistics slide"""
        slide = self.add_content_slide("Statystyki tygodniowe", 9)
        stats = self.weather_data['stats']
        
        metrics = [
            ('Średnia temp', f"{stats['avg_temp']}°C", True),
            ('Maksymalna', f"{stats['max_temp']}°C", False),
            ('Minimalna', f"{stats['min_temp']}°C", False)
        ]
        
        x = 0.5
        for label, value, is_red in metrics:
            # Box
            box = slide.shapes.add_shape(1, Inches(x), Inches(2.5), Inches(2.8), Inches(1.8))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(255, 230, 230) if is_red else self.light_gray
            box.line.color.rgb = self.gray
            
            # Label
            lbl = slide.shapes.add_textbox(Inches(x+0.2), Inches(2.7), Inches(2.4), Inches(0.4))
            lbl.text_frame.text = label
            lbl.text_frame.paragraphs[0].font.size = Pt(12)
            lbl.text_frame.paragraphs[0].font.color.rgb = self.gray
            lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Value
            val = slide.shapes.add_textbox(Inches(x+0.2), Inches(3.2), Inches(2.4), Inches(0.8))
            val.text_frame.text = value
            val.text_frame.paragraphs[0].font.size = Pt(42)
            val.text_frame.paragraphs[0].font.color.rgb = self.red if is_red else self.black
            val.text_frame.paragraphs[0].font.bold = True
            val.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            x += 3
    
    def add_summary(self):
        """Summary slide"""
        slide = self.add_content_slide("Podsumowanie – Warszawa", 16)
        
        summaries = [
            ('Stabilne warunki', f"Średnia temperatura {self.weather_data['stats']['avg_temp']}°C z niewielkimi wahaniami."),
            ('Umiarkowane opady', f"Przewidywane {self.weather_data['stats']['rainy_days']} dni z opadami."),
            ('Komfort termiczny', f"Temperatura odczuwalna {self.weather_data['current']['temp']}°C zapewnia komfort.")
        ]
        
        y = 2
        for title, text in summaries:
            # Title
            t = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(8), Inches(0.4))
            t.text_frame.text = title
            t.text_frame.paragraphs[0].font.size = Pt(18)
            t.text_frame.paragraphs[0].font.color.rgb = self.red
            t.text_frame.paragraphs[0].font.bold = True
            
            # Text
            txt = slide.shapes.add_textbox(Inches(0.5), Inches(y+0.4), Inches(8), Inches(0.6))
            txt.text_frame.text = text
            txt.text_frame.paragraphs[0].font.size = Pt(12)
            txt.text_frame.paragraphs[0].font.color.rgb = self.black
            txt.text_frame.word_wrap = True
            
            y += 1.5
    
    def add_final(self):
        """Final logo slide"""
        slide = self.add_blank_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.white
        
        logo = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(4), Inches(1.5))
        tf = logo.text_frame
        tf.text = "WEATHER\nSERVICE"
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in tf.paragraphs:
            p.font.size = Pt(60)
            p.font.color.rgb = self.red
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
    
    def generate(self):
        """Generate complete presentation"""
        print("Generating weather presentation...")
        
        self.add_title_slide()
        print("  ✓ Title slide")
        
        self.add_toc_slide()
        print("  ✓ Table of contents")
        
        self.add_section_divider(1, "Kontekst i dane\npodstawowe")
        print("  ✓ Section 1 divider")
        
        self.add_current_weather()
        print("  ✓ Current weather")
        
        self.add_section_divider(2, "Prognoza\ntygodniowa")
        print("  ✓ Section 2 divider")
        
        # Placeholder slide
        slide = self.add_blank_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.white
        self.add_footer(slide, 6)
        
        self.add_weekly_forecast()
        print("  ✓ Weekly forecast")
        
        self.add_section_divider(3, "Analiza\ntemperatur")
        print("  ✓ Section 3 divider")
        
        self.add_stats()
        print("  ✓ Statistics")
        
        self.add_section_divider(4, "Opady\ni wiatr")
        print("  ✓ Section 4 divider")
        
        # Placeholder slides
        for i in range(11, 15):
            slide = self.add_blank_slide()
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = self.white
            self.add_footer(slide, i)
        
        self.add_section_divider(5, "Podsumowanie")
        print("  ✓ Section 5 divider")
        
        self.add_summary()
        print("  ✓ Summary")
        
        self.add_final()
        print("  ✓ Final slide")
        
        # Save
        output_path = os.path.join(self.output_dir, 'Prognoza_Pogody_Warszawa.pptx')
        os.makedirs(self.output_dir, exist_ok=True)
        self.prs.save(output_path)
        
        print(f"\n✓ Presentation saved: {output_path}")
        return output_path

def main():
    try:
        gen = WeatherPresentation(output_dir='test_02')
        gen.generate()
        return 0
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        return 1

if __name__ == '__main__':
    sys.exit(main())
