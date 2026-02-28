#!/usr/bin/env python3
"""
STEP 1 of 3: Extract PPTX slide data and render screenshots.

Saves intermediates to .tmp/{stem}/ for Claude to analyse (Step 2),
then build_document.py (Step 3) assembles the final PDF.

Usage:
    python3 execution/convert_for_ai.py [input_folder]

Default input_folder: file_converter/Input
Output: .tmp/{stem}/slides/*.png  +  .tmp/{stem}/extracted.json
"""

import sys
import json
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("ERROR: pip3 install python-pptx")
    sys.exit(1)


BASE = Path(__file__).parent.parent  # project root


# ─── Text extraction ───────────────────────────────────────────────────────────

def extract_text_from_shapes(shapes) -> list[str]:
    """Recursively extract non-empty text. Skips footer/slide-number/date placeholders."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    SKIP_PH_TYPES = {PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER, PP_PLACEHOLDER.DATE}
    texts = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            texts.extend(extract_text_from_shapes(shape.shapes))
        else:
            try:
                ph = shape.placeholder_format
                if ph and ph.type in SKIP_PH_TYPES:
                    continue
            except Exception:
                pass
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)
    return texts


# ─── PPTX → .tmp/ ─────────────────────────────────────────────────────────────

def process_pptx(pptx_path: Path, tmp_root: Path) -> Path:
    """
    Renders slides as PNG and extracts text to .tmp/{stem}/.
    Returns the tmp dir path.
    """
    prs = Presentation(str(pptx_path))
    stem = pptx_path.stem
    tmp_dir = tmp_root / stem
    slides_dir = tmp_dir / "slides"

    # Clean stale PNGs
    if slides_dir.exists():
        for old in slides_dir.glob("slide_*.png"):
            old.unlink()
    slides_dir.mkdir(parents=True, exist_ok=True)

    # Render slides via Spire.Presentation (full visual rendering — colors, shapes, diagrams)
    print(f"  Rendering {len(prs.slides)} slides...")
    from spire.presentation import Presentation as SpirePrs
    spire = SpirePrs()
    spire.LoadFromFile(str(pptx_path))
    for i, spire_slide in enumerate(spire.Slides):
        img = spire_slide.SaveAsImage()
        img.Save(str(slides_dir / f"slide_{i + 1:02d}.png"))
        img.Dispose()
    spire.Dispose()

    # Extract text per slide
    slides_data = []
    for i, slide in enumerate(prs.slides):
        title = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text_frame.text.strip()

        all_texts = extract_text_from_shapes(slide.shapes)
        # Deduplicate: remove title from body texts (it's already in the heading)
        texts = [t for t in all_texts if t != title]

        # Speaker notes
        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        slides_data.append({
            "slide": i + 1,
            "title": title,
            "texts": texts,
            "notes": notes,
        })

    extracted_path = tmp_dir / "extracted.json"
    extracted_path.write_text(json.dumps(slides_data, ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"  Saved {len(prs.slides)} PNGs → {slides_dir}")
    print(f"  Saved extracted text → {extracted_path}")
    return tmp_dir


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    input_dir = Path(sys.argv[1]) if len(sys.argv) > 1 else BASE / "file_converter" / "Input"
    input_dir = input_dir.resolve()

    if not input_dir.exists():
        print(f"ERROR: Input folder not found: {input_dir}")
        sys.exit(1)

    tmp_root = BASE / ".tmp"

    pptx_files = list(input_dir.glob("*.pptx")) + list(input_dir.glob("*.PPTX"))

    if not pptx_files:
        print(f"No PPTX files found in: {input_dir}")
        sys.exit(0)

    print(f"Found {len(pptx_files)} PPTX\n")

    tmp_dirs = []

    for pptx_path in pptx_files:
        print(f"[PPTX] {pptx_path.name}")
        try:
            tmp_dir = process_pptx(pptx_path, tmp_root)
            tmp_dirs.append(tmp_dir)
        except Exception as e:
            print(f"  ERROR: {e}")

    if tmp_dirs:
        print("\n" + "─" * 60)
        print("STEP 1 done. Next steps:")
        print()
        print("STEP 2: Claude reviews slides and writes descriptions.")
        print("  For each folder below, Claude reads slide PNGs and")
        print("  generates .tmp/{stem}/descriptions.json")
        print()
        for d in tmp_dirs:
            print(f"  {d}")
        print()
        print("STEP 3: Build final PDF:")
        for d in tmp_dirs:
            print(f"  python3 execution/build_document.py \"{d}\"")


if __name__ == "__main__":
    main()
